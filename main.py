import requests
import os
import openai
from tqdm import tqdm
import json
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor


topic = input("Topic : ")
slide = False
while not slide:
    try:
        slide = int(input("No. of slide : "))
    except:
        print("Enter a number")

reply1 = ""
def getBasicPPT(topic,slide):
    global reply1
    openai.api_key = os.environ['api_key']
    messages = [
        {"role": "system", "content": "You are a assistant."},
    ]
    msgs = [
        os.environ['msg1'].replace("{topic}",topic).replace("{slide}",str(slide)),
        os.environ['msg2']
    ]
    try:
        for message in tqdm(msgs):
            messages.append(
                {"role": "user", "content": message},
            )
            chat = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", messages = messages
            )
            reply = chat.choices[0].message.content
            messages.append({"role": "assistant", "content": reply})
    except:
        print(reply)
        exit()
    reply1 = reply
    reply = str(reply)
    reply = reply.removeprefix(reply.split("{")[0])
    reply = reply.removesuffix(reply.split("}")[-1])
    with open("reply.json", "w") as f:
        f.write(reply)


def createPPT():
    try:
        with open('reply.json') as f:
            data = json.load(f)
        prs = Presentation()
        for page in data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(216, 223, 219)
            slide.colour = "blue"
            title = slide.shapes.title
            try:
                title.text = data[page]['title']
            except:
                title.text = data[page]['Title']
            textframe = title.text_frame
            paragraph = textframe.paragraphs[0]
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.underline = True
            points = ""
            try:
                content = data[page]['content']
            except:
                content = data[page]['Content']
            for point in content:
                points+=str(point).removeprefix("-")+"\n"
            subtitle = slide.placeholders[1]
            subtitle.text = points
            subtitle.width = Pt(400)
            subtitle.height = Pt(350)
            textframe = subtitle.text_frame
            for paragraph in textframe.paragraphs:
                paragraph.font.size = Pt(17)
        prs.save('result.pptx')
    except:
        print("\n\nERROR  :")
        print(reply1)
        exit()

def getDownloadLink():
    url = 'https://file.io'
    file_path = './result.pptx'
    with open(file_path, 'rb') as file:
        response = requests.post(url, files={'file': file})
    json_response = response.json()
    download_link = json_response['link']
    print("\n\nDownload Link : ",download_link)



getBasicPPT(topic,slide)
createPPT()
getDownloadLink()


